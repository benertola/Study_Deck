import io
import os


def parse_bytes(filename: str, data: bytes) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return _parse_pdf(data)
    elif ext == ".pptx":
        return _parse_pptx(data)
    elif ext == ".docx":
        return _parse_docx(data)
    elif ext == ".txt":
        return data.decode("utf-8", errors="ignore")
    return ""


def _parse_pdf(data: bytes) -> str:
    import pdfplumber

    pages = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n".join(pages)


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _parse_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
